import os
import glob
import shutil
from pptx import Presentation
from pptx.util import Inches

base_dir = '/Volumes/legion-laptop.local/code/VietIS/Vibe Coding/Tài liệu miễn phí tặng khách'
logo_path = os.path.join(base_dir, 'logo.png')

# Loop through folders 01 to 10
for i in range(1, 11):
    folder_name = f"{i:02d}"
    folder_path = os.path.join(base_dir, folder_name)
    
    if not os.path.isdir(folder_path):
        continue
        
    # Find pptx in the folder
    pptx_files = glob.glob(os.path.join(folder_path, '*.pptx'))
    if not pptx_files:
        print(f"No pptx found in {folder_name}")
        continue
        
    pptx_path = pptx_files[0]
    
    # Create a new presentation
    new_prs = Presentation()
    new_prs.slide_width = Inches(13.333)
    new_prs.slide_height = Inches(7.5)
    
    # Get all images in the folder
    images = sorted(glob.glob(os.path.join(folder_path, '*.jpg')) + glob.glob(os.path.join(folder_path, '*.png')))
    images = [img for img in images if os.path.basename(img) != 'logo.png']
    
    layout = new_prs.slide_layouts[6] if len(new_prs.slide_layouts) > 6 else new_prs.slide_layouts[0]
    
    for img in images:
        slide = new_prs.slides.add_slide(layout)
        # remove default shapes
        for shape in list(slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
            
        # Add background image
        slide.shapes.add_picture(img, 0, 0, width=new_prs.slide_width, height=new_prs.slide_height)
        
        # Add logo
        try:
            slide.shapes.add_picture(logo_path, Inches(0.05), Inches(0.05), width=Inches(2.5))
        except Exception as e:
            print(f"Error adding logo to {img}: {e}")
            
    # Save the updated presentation
    new_prs.save(pptx_path)
    print(f"Processed {folder_name}: Added {len(images)} slides to {os.path.basename(pptx_path)}")
    
    # Move the pptx file to the base directory
    dest_path = os.path.join(base_dir, os.path.basename(pptx_path))
    shutil.move(pptx_path, dest_path)
    print(f"Moved {os.path.basename(pptx_path)} to base directory.")

print("All tasks completed successfully!")
