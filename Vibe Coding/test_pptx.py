import os
import glob
from pptx import Presentation
from pptx.util import Inches

base_dir = '/Volumes/legion-laptop.local/code/VietIS/Vibe Coding/Tài liệu miễn phí tặng khách'
folder_01 = os.path.join(base_dir, '01')
logo_path = os.path.join(base_dir, 'logo.png')

# Find pptx in 01
pptx_files = glob.glob(os.path.join(folder_01, '*.pptx'))
if not pptx_files:
    print("No pptx found in 01")
    exit()

pptx_path = pptx_files[0]

# Create a new presentation to avoid zipfile duplicates
new_prs = Presentation()

# Force 16:9 slide size (widescreen default for modern PPT)
new_prs.slide_width = Inches(13.333)
new_prs.slide_height = Inches(7.5)

# Get images in 01
images = sorted(glob.glob(os.path.join(folder_01, '*.jpg')) + glob.glob(os.path.join(folder_01, '*.png')))
# Just take first 2 for testing
images = [img for img in images if os.path.basename(img) != 'logo.png'][:2]

layout = new_prs.slide_layouts[6] if len(new_prs.slide_layouts) > 6 else new_prs.slide_layouts[0]

for img in images:
    slide = new_prs.slides.add_slide(layout)
    # delete any shapes that might be in the layout
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Add background image full screen
    slide.shapes.add_picture(img, 0, 0, width=new_prs.slide_width, height=new_prs.slide_height)
    
    # Add logo top left: closer to edge (0.05 inches) and larger (2.5 inches width)
    try:
        slide.shapes.add_picture(logo_path, Inches(0.05), Inches(0.05), width=Inches(2.5))
    except Exception as e:
        print(f"Error adding logo: {e}")

# Save replacing the old file
new_prs.save(pptx_path)
print(f"Saved {pptx_path} with 2 slides (16:9, larger logo).")
