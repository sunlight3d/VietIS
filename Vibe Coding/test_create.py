from pptx import Presentation

try:
    prs = Presentation('VIETIS.EDU_TRN_Template.pptx')

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    
    for shape in slide1.placeholders:
        if shape.placeholder_format.type == 3: # CENTER_TITLE
            shape.text = "TỔNG HỢP CÔNG CỤ AI & CÔNG THỨC RCESF"
        elif shape.placeholder_format.type == 4: # SUBTITLE
            shape.text = "Bí quyết tăng hiệu suất làm việc với AI Tools\n(Dựa trên tài liệu VisionEdu)"

    # Slide 2: Content Slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Try to set title if it exists, otherwise add a textbox
    if slide2.shapes.title:
        slide2.shapes.title.text = "Công thức RCESF - Prompt chuẩn chỉnh"
    else:
        # Fallback if no title shape
        from pptx.util import Pt, Inches
        txBox = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = "Công thức RCESF - Prompt chuẩn chỉnh"
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True

    for shape in slide2.placeholders:
        if shape.placeholder_format.type == 2: # BODY
            tf = shape.text_frame
            tf.text = "R (Role): Giao đúng vai cho AI (Ví dụ: Bạn là chuyên gia Marketing...)"
            
            p = tf.add_paragraph()
            p.text = "C (Context): Đặt đúng bối cảnh (Ví dụ: Tôi cần viết email báo cáo...)"
            
            p = tf.add_paragraph()
            p.text = "E (Example): Cho AI ví dụ hoặc mẫu tham chiếu"
            
            p = tf.add_paragraph()
            p.text = "S (Style): Quy định giọng điệu (Ví dụ: Ngắn gọn, chuyên nghiệp)"
            
            p = tf.add_paragraph()
            p.text = "F (Format): Gắn định dạng đầu ra (Ví dụ: 5 gạch đầu dòng)"

    prs.save('Draft_AI_Presentation.pptx')
    print("Successfully created Draft_AI_Presentation.pptx")
except Exception as e:
    print(f"Error: {e}")
